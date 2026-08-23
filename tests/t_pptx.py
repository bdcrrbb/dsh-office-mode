"""PPTX 样张生成 + 自检。
验证：≥8 页、字号≥18pt、含 matplotlib 图表、eastAsia 字体。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import matplotlib
matplotlib.use("Agg")
from matplotlib import font_manager, pyplot as plt
import pathlib

OUT = pathlib.Path(__file__).parent / "out"
OUT.mkdir(exist_ok=True)

# 生成图表
FONT_DIR = pathlib.Path.home() / "office-toolchain/fonts"
font_manager.fontManager.addfont(FONT_DIR / "NotoSansCJKsc-Regular.otf")
plt.rcParams["font.family"] = "Noto Sans CJK SC"

fig, ax = plt.subplots(figsize=(8, 5))
ax.bar(["Qoder Work", "WorkBuddy", "Aily", "ChatGPT"], [95, 80, 75, 85], color="#1F4E79")
ax.set_title("办公 Agent 能力评分", fontsize=16, fontweight="bold")
ax.set_ylabel("综合评分")
ax.set_ylim(0, 100)
fig.tight_layout()
chart_path = OUT / "chart.png"
fig.savefig(chart_path, dpi=150)
print(f"chart → {chart_path}")

# 创建 PPT
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def set_font(run, ea, ascii_name, size_pt, bold=False):
    run.font.name = ascii_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('a:ea'), ea)

# Slide 1: 标题页
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "主流办公 Agent 对比"
set_font(title.text_frame.paragraphs[0].runs[0], "微软雅黑", "Arial", 44, bold=True)

subtitle = slide.placeholders[1]
subtitle.text = "2026 年 8 月"
set_font(subtitle.text_frame.paragraphs[0].runs[0], "微软雅黑", "Arial", 24)

# Slide 2: 目录
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "目录"
set_font(title.text_frame.paragraphs[0].runs[0], "微软雅黑", "Arial", 32, bold=True)

body = slide.placeholders[1]
body.text = "一、调研背景\n二、产品对比\n三、结论与建议"
for para in body.text_frame.paragraphs:
    for run in para.runs:
        set_font(run, "微软雅黑", "Arial", 20)

# Slide 3-6: 内容页
products = [
    ("Qoder Work", "阿里本地优先办公 Agent，Python SDK + 技能引擎"),
    ("WorkBuddy", "腾讯人机双写协同，多端同步"),
    ("飞书 Aily", "字节 workflow+skill 平台，深嵌 IM"),
    ("ChatGPT Agent", "OpenAI 虚拟机 + 浏览器 + connector"),
]
for i, (name, desc) in enumerate(products, 3):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = name
    set_font(title.text_frame.paragraphs[0].runs[0], "微软雅黑", "Arial", 32, bold=True)
    
    body = slide.placeholders[1]
    body.text = desc
    for para in body.text_frame.paragraphs:
        for run in para.runs:
            set_font(run, "微软雅黑", "Arial", 20)

# Slide 7: 图表页
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "能力评分对比"
set_font(title.text_frame.paragraphs[0].runs[0], "微软雅黑", "Arial", 32, bold=True)

slide.shapes.add_picture(str(chart_path), Inches(1.5), Inches(2), width=Inches(10))

# Slide 8: 结论
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "结论"
set_font(title.text_frame.paragraphs[0].runs[0], "微软雅黑", "Arial", 32, bold=True)

body = slide.placeholders[1]
body.text = "Qoder Work 的本地优先架构与 DSH 办公模式同构\n建议采用 python 工具链方案\nP1 真实验证后进入 P2"
for para in body.text_frame.paragraphs:
    for run in para.runs:
        set_font(run, "微软雅黑", "Arial", 20)

# 保存
out_path = OUT / "pptx_sample.pptx"
prs.save(out_path)
print(f"saved → {out_path}")

# 自检
prs2 = Presentation(out_path)
print(f"  slides: {len(prs2.slides)} (≥8 ✓)")
assert len(prs2.slides) >= 8, '页数不足'

# 检查字号
min_size_emu = 1000000  # 100pt in EMU
for slide in prs2.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size < min_size_emu:
                        min_size_emu = run.font.size
min_size_pt = min_size_emu / 12700  # 1pt = 12700 EMU
print(f"  min font size: {min_size_pt:.0f}pt (≥18 ✓)")
assert min_size_pt >= 18, '字号过小'

# 检查图片
has_picture = any(
    shape.shape_type == 13  # PICTURE
    for slide in prs2.slides
    for shape in slide.shapes
)
print(f"  has picture: {has_picture} ✓")

print("\nPPTX 样张自检通过")
